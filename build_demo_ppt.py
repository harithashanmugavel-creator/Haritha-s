"""
Build a demo PowerPoint deck for OptiGuard AI - Marketing Analytics AI app.

Run: python3 build_demo_ppt.py
Output: OptiGuard_AI_Demo.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# -----------------------------------------------------------------------------
# Theme
# -----------------------------------------------------------------------------
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
DEEP = RGBColor(0x16, 0x21, 0x3E)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PINK = RGBColor(0xF4, 0x72, 0xB6)
GREEN = RGBColor(0x28, 0xA7, 0x45)
AMBER = RGBColor(0xFF, 0xA7, 0x26)
RED = RGBColor(0xDC, 0x35, 0x45)
LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
GREY = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=NAVY):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=LIGHT,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.text = ""
    lines = text.split("\n") if isinstance(text, str) else [str(text)]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=LIGHT, bullet="•  "):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet}{item}"
        run.font.size = Pt(size)
        run.font.name = "Calibri"
        run.font.color.rgb = color
        p.space_after = Pt(4)
    return tb


def add_card(slide, x, y, w, h, *, fill=DEEP, border=CYAN, border_w=1.25):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_w)
    card.shadow.inherit = False
    return card


def add_pill(slide, x, y, w, h, text, *, fill=CYAN, color=NAVY, size=11, bold=True):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    pill.shadow.inherit = False
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return pill


def slide_header(slide, title, subtitle=None, accent=CYAN):
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False

    add_text(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7),
             title, size=30, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4),
                 subtitle, size=14, color=GREY)

    # Footer
    add_text(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
             "OptiGuard AI  •  Marketing Operations Control Tower", size=9, color=GREY)


# -----------------------------------------------------------------------------
# Slide 1 – Title
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

# Decorative gradient block
deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), SW, Inches(0.04))
deco.fill.solid(); deco.fill.fore_color.rgb = CYAN; deco.line.fill.background()

add_text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(1.0),
         "🛡️  OptiGuard AI", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(2.55), Inches(12), Inches(0.7),
         "Campaign Intelligence & ROI Optimization Copilot", size=24, color=CYAN)
add_text(s, Inches(0.7), Inches(3.2), Inches(12), Inches(0.6),
         "Predictive Marketing Operations Intelligence for Banking",
         size=16, color=GREY)

# Three highlight pills
pills = [
    ("Predict failures BEFORE launch", PURPLE),
    ("AI-powered risk scoring", CYAN),
    ("Historical learning via RAG", PINK),
]
x = Inches(0.7)
for label, col in pills:
    add_pill(s, x, Inches(4.2), Inches(3.6), Inches(0.45), label,
             fill=col, color=WHITE, size=12)
    x += Inches(3.85)

add_text(s, Inches(0.7), Inches(5.4), Inches(12), Inches(0.5),
         "Demo Walkthrough  |  Architecture  |  Outputs",
         size=18, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.4),
         "Powered by Azure OpenAI (GPT) + Embeddings + FAISS + Streamlit",
         size=12, color=GREY)


# -----------------------------------------------------------------------------
# Slide 2 – The Problem
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "The Problem", "Why marketing operations teams need predictive intelligence")

problems = [
    ("📊", "Fragmented Intelligence", "Campaign data scattered across SMS, Email, WhatsApp, Push & Call Center"),
    ("🔁", "Repeated Targeting", "Same customers hit again & again — fatigue & opt-outs spike"),
    ("⏱️", "Delayed ROI Analysis", "Insights arrive AFTER campaigns fail — pure post-mortem"),
    ("⚖️", "Manual Compliance", "Risky wording & regulatory issues caught too late"),
    ("📚", "Lost Learnings", "Past wins/failures buried in spreadsheets — no retrieval"),
    ("💸", "Cost Impact", "20–30 analyst hours / campaign  →  ₹7.2 Cr / year"),
]
for i, (icon, title, desc) in enumerate(problems):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.25)
    y = Inches(1.6 + row * 2.55)
    add_card(s, x, y, Inches(4.0), Inches(2.3), border=RED)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(0.6), Inches(0.6), icon, size=28)
    add_text(s, x + Inches(0.95), y + Inches(0.25), Inches(3.0), Inches(0.5),
             title, size=16, bold=True, color=CYAN)
    add_text(s, x + Inches(0.2), y + Inches(0.95), Inches(3.7), Inches(1.3),
             desc, size=12, color=LIGHT)


# -----------------------------------------------------------------------------
# Slide 3 – Solution Overview
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "The Solution",
             "OptiGuard AI — a predictive Campaign Risk Intelligence platform")

add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
         "From REACTIVE dashboards  →  to  PREDICTIVE intelligence",
         size=20, bold=True, color=PINK)

caps = [
    ("🔮", "Predict Risk Before Launch",
     "Fatigue • Opt-out • Complaint • Trust • Channel safety"),
    ("🧠", "AI Reasoning + Evidence",
     "GPT explains WHY scores were given, backed by historical data"),
    ("🔍", "Semantic Search (RAG)",
     "Find similar past campaigns instantly with FAISS + embeddings"),
    ("🎯", "Audience Playbooks",
     "Best timing, channel, segments, messaging recommendations"),
    ("⚡", "Campaign Optimizer",
     "AI rewrites the campaign to lower risk & raise predicted ROI"),
    ("📈", "Executive Summary",
     "Auto-generated narrative for leadership with impact estimates"),
]
for i, (icon, title, desc) in enumerate(caps):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.25)
    y = Inches(2.2 + row * 2.35)
    add_card(s, x, y, Inches(4.0), Inches(2.15), border=CYAN)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(0.6), Inches(0.6), icon, size=24)
    add_text(s, x + Inches(0.85), y + Inches(0.22), Inches(3.0), Inches(0.5),
             title, size=14, bold=True, color=CYAN)
    add_text(s, x + Inches(0.2), y + Inches(0.85), Inches(3.7), Inches(1.2),
             desc, size=11, color=LIGHT)


# -----------------------------------------------------------------------------
# Slide 4 – Architecture
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Architecture",
             "Streamlit  →  AI/RAG layer  →  Risk engine  →  Data foundation")

# 4 horizontal layers
layers = [
    ("PRESENTATION LAYER", "control_tower_app.py  •  risk_intelligence_app.py  •  app.py (Streamlit UIs)", CYAN),
    ("AI / LLM ORCHESTRATION", "Azure OpenAI GPT (gpt-5.4-gamma)  •  text-embedding-3-small  •  Prompt templates  •  JSON contracts", PURPLE),
    ("INTELLIGENCE ENGINES", "CampaignRiskPredictor  •  SimilarCampaignFinder (FAISS)  •  Audience Playbook  •  Optimizer  •  Exec Summary", PINK),
    ("DATA FOUNDATION", "campaign_history (10K)  •  campaign_content (20K)  •  audience_exposure (50K)  •  feedback (20K)  •  ROI (5K)  •  QA (5K)", AMBER),
]
y = Inches(1.55)
for label, contents, col in layers:
    add_card(s, Inches(0.5), y, Inches(12.3), Inches(1.15), border=col)
    add_text(s, Inches(0.7), y + Inches(0.1), Inches(11.8), Inches(0.4),
             label, size=14, bold=True, color=col)
    add_text(s, Inches(0.7), y + Inches(0.55), Inches(11.8), Inches(0.55),
             contents, size=12, color=LIGHT)
    # arrow between layers
    if y < Inches(5.5):
        ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.55), y + Inches(1.18),
                                Inches(0.25), Inches(0.18))
        ar.fill.solid(); ar.fill.fore_color.rgb = col; ar.line.fill.background()
    y += Inches(1.36)


# -----------------------------------------------------------------------------
# Slide 5 – Tech Stack & Data Flow
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Tech Stack & End-to-End Data Flow",
             "How a campaign idea becomes a predicted risk score")

# Left column: stack
add_text(s, Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.4),
         "Technology Stack", size=18, bold=True, color=CYAN)
add_card(s, Inches(0.5), Inches(1.95), Inches(5.5), Inches(4.8))
stack = [
    "UI / UX:  Streamlit (dark glassmorphism theme)",
    "LLM:  Azure OpenAI GPT (gpt-5.4-gamma)",
    "Embeddings:  text-embedding-3-small (1536-dim)",
    "Vector DB:  FAISS (IndexFlatIP) — local index",
    "Data:  Pandas on 6 banking-domain CSVs",
    "Caching:  @st.cache_resource for indices/models",
    "Config:  python-dotenv (.env for keys)",
    "Language:  Python 3  •  Logging built-in",
    "RAG pattern:  semantic retrieval → prompt → JSON",
    "Pickled metadata:  campaign_metadata.pkl",
]
add_bullets(s, Inches(0.7), Inches(2.1), Inches(5.2), Inches(4.5),
            stack, size=12)

# Right column: flow
add_text(s, Inches(6.4), Inches(1.5), Inches(6.4), Inches(0.4),
         "Request Flow", size=18, bold=True, color=PINK)

steps = [
    ("1", "User enters campaign idea + audience + channel + budget", CYAN),
    ("2", "App loads CSVs → builds historical patterns (segment / channel / region)", CYAN),
    ("3", "FAISS retrieves top-K similar past campaigns via embeddings", PURPLE),
    ("4", "Risk engine computes Fatigue, Opt-out, Complaint, Trust, Channel-safety, ROI", PURPLE),
    ("5", "GPT receives scores + history → returns reasoning, insights, recommendations (JSON)", PINK),
    ("6", "Streamlit renders Health Score, factors, playbook, optimizer & exec summary", PINK),
]
y = Inches(1.95)
for num, txt, col in steps:
    add_card(s, Inches(6.4), y, Inches(6.4), Inches(0.72), border=col, border_w=1.0)
    cir = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.55), y + Inches(0.16),
                             Inches(0.4), Inches(0.4))
    cir.fill.solid(); cir.fill.fore_color.rgb = col; cir.line.fill.background()
    tf = cir.text_frame; tf.margin_left = tf.margin_right = Inches(0.0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = NAVY
    add_text(s, Inches(7.05), y + Inches(0.18), Inches(5.3), Inches(0.5),
             txt, size=11, color=LIGHT)
    y += Inches(0.78)


# -----------------------------------------------------------------------------
# Slide 6 – Datasets
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Data Foundation",
             "6 synthetic banking marketing datasets — 110,000 rows total")

datasets = [
    ("campaign_history.csv", "10,000",
     "Campaign metrics: spend, CTR, conversions, ROI, status, objective", CYAN),
    ("campaign_content.csv", "20,000",
     "SMS / Email / WhatsApp body, tone, CTA, compliance risk, A/B variant", PURPLE),
    ("audience_exposure.csv", "50,000",
     "7-day & 30-day exposure counts, fatigue & churn scores, opt-out risk", PINK),
    ("campaign_feedback.csv", "20,000",
     "Customer complaints, sentiment, severity, escalations, repeat issues", AMBER),
    ("campaign_roi_metrics.csv", "5,000",
     "ROAS, CAC, LTV, profit, attribution model, performance vs. benchmark", GREEN),
    ("campaign_qa_results.csv", "5,000",
     "Tone consistency, CTA alignment, brand & regulatory compliance scores", CYAN),
]
for i, (name, rows, desc, col) in enumerate(datasets):
    cidx = i % 2
    ridx = i // 2
    x = Inches(0.5 + cidx * 6.4)
    y = Inches(1.6 + ridx * 1.75)
    add_card(s, x, y, Inches(6.2), Inches(1.55), border=col)
    add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(4.5), Inches(0.45),
             name, size=14, bold=True, color=col)
    add_pill(s, x + Inches(4.9), y + Inches(0.15), Inches(1.1), Inches(0.35),
             rows + " rows", fill=col, color=NAVY, size=10)
    add_text(s, x + Inches(0.2), y + Inches(0.65), Inches(5.8), Inches(0.85),
             desc, size=11, color=LIGHT)


# -----------------------------------------------------------------------------
# Slide 7 – How the App Works (User Journey)
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "How the App Works",
             "End-to-end user journey inside the Control Tower")

journey = [
    ("①  Configure", "User fills the sidebar:\n• Campaign idea\n• Audience segment\n• Primary + secondary channels\n• Product, region, budget\n• Launch & end dates", CYAN),
    ("②  Predict",  "Click PREDICT RISKS:\n• Risk engine scores 5 dimensions\n• FAISS retrieves similar campaigns\n• GPT generates insights + reasoning\n• Health score is calculated", PURPLE),
    ("③  Review",   "Results panel shows:\n• Big Campaign Health Score\n• Subscores & contributing factors\n• Audience playbook\n• Similar historical campaigns", PINK),
    ("④  Act",      "User can:\n• Click OPTIMIZE CAMPAIGN to rewrite\n• Chat contextually with each section\n• Export executive summary\n• Decide PROCEED / CAUTION / RECONSIDER", AMBER),
]
for i, (title, body, col) in enumerate(journey):
    x = Inches(0.5 + i * 3.18)
    y = Inches(1.6)
    add_card(s, x, y, Inches(3.0), Inches(5.1), border=col)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(2.7), Inches(0.5),
             title, size=18, bold=True, color=col)
    add_text(s, x + Inches(0.2), y + Inches(0.8), Inches(2.7), Inches(4.2),
             body, size=11, color=LIGHT)


# -----------------------------------------------------------------------------
# Slide 8 – Output 1: Campaign Health Score
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Output 1 — Campaign Health Score",
             "The single hero metric: 0–100 grade with reasoning")

# big score circle
circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(1.7),
                            Inches(3.6), Inches(3.6))
circle.fill.solid(); circle.fill.fore_color.rgb = DEEP
circle.line.color.rgb = CYAN; circle.line.width = Pt(3); circle.shadow.inherit = False
add_text(s, Inches(0.8), Inches(2.3), Inches(3.6), Inches(1.2),
         "78", size=88, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(3.7), Inches(3.6), Inches(0.5),
         "Grade:  B+", size=18, bold=True, color=PINK, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(4.2), Inches(3.6), Inches(0.6),
         "Healthy — Proceed with monitoring", size=12, color=GREY,
         align=PP_ALIGN.CENTER)

# subscores
add_text(s, Inches(4.8), Inches(1.7), Inches(8.0), Inches(0.45),
         "Health Subscores", size=18, bold=True, color=CYAN)
subs = [
    ("💰  ROI Health",      82, GREEN),
    ("🤝  Trust Health",    74, CYAN),
    ("😊  Engagement Health", 71, PURPLE),
    ("🛡️  Compliance Health", 88, GREEN),
    ("📡  Channel Health",  68, AMBER),
]
y = Inches(2.25)
for label, val, col in subs:
    add_text(s, Inches(4.85), y, Inches(3.2), Inches(0.4),
             label, size=12, bold=True, color=LIGHT)
    # bar background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.05), y + Inches(0.08),
                            Inches(4.0), Inches(0.22))
    bg.fill.solid(); bg.fill.fore_color.rgb = DEEP
    bg.line.color.rgb = GREY; bg.line.width = Pt(0.5); bg.shadow.inherit = False
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.05), y + Inches(0.08),
                             Inches(4.0 * val / 100), Inches(0.22))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    bar.shadow.inherit = False
    add_text(s, Inches(12.15), y, Inches(0.7), Inches(0.4),
             f"{val}", size=12, bold=True, color=col)
    y += Inches(0.55)

# Interpretation
add_card(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.3), border=PURPLE)
add_text(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.4),
         "AI Interpretation", size=13, bold=True, color=PURPLE)
add_text(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.85),
         "“Strong predicted ROI and clean compliance posture, but Channel Health is "
         "trending low because Salaried Millennials have already been hit 4× via SMS in the "
         "past 30 days. Switch primary channel to Email or shift launch by 7 days.”",
         size=12, color=LIGHT)


# -----------------------------------------------------------------------------
# Slide 9 – Output 2: Risk Scores
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Output 2 — Predictive Risk Scores",
             "5 risk dimensions, each with level + AI description")

risks = [
    ("Fatigue Risk",       62, "MEDIUM",  AMBER, "Audience hit 4× in last 30 days via SMS"),
    ("Opt-Out Probability", 38, "MEDIUM",  AMBER, "≈ 38% chance of unsubscribes"),
    ("Complaint Likelihood", 22, "LOW",    GREEN, "Historical complaint rate well below threshold"),
    ("Trust Risk",         41, "MEDIUM",  AMBER, "Combined fatigue + opt-out lifting trust risk"),
    ("Channel Safety",     71, "GOOD",    GREEN, "Email is among the safest channels for HNI"),
]
for i, (label, val, lvl, col, desc) in enumerate(risks):
    cidx = i % 3
    ridx = i // 3
    x = Inches(0.5 + cidx * 4.25)
    y = Inches(1.6 + ridx * 2.45)
    add_card(s, x, y, Inches(4.0), Inches(2.25), border=col)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(3.7), Inches(0.4),
             label, size=14, bold=True, color=col)
    add_text(s, x + Inches(0.2), y + Inches(0.55), Inches(2.4), Inches(0.9),
             f"{val}%", size=36, bold=True, color=col)
    add_pill(s, x + Inches(2.7), y + Inches(0.7), Inches(1.1), Inches(0.4),
             lvl, fill=col, color=NAVY, size=11)
    add_text(s, x + Inches(0.2), y + Inches(1.55), Inches(3.7), Inches(0.65),
             desc, size=11, color=LIGHT)

# Overall pill
add_card(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.6), border=CYAN)
add_text(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.5),
         "Overall Risk: 47%   |   Level: MEDIUM   |   Launch Recommendation: ⚡ CAUTION",
         size=14, bold=True, color=CYAN)


# -----------------------------------------------------------------------------
# Slide 10 – Output 3: Risk Factor Breakdown + Similar Campaigns
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Output 3 — Why?  Factor Breakdown + Historical Evidence",
             "AI explains the score using contributing factors and similar past campaigns")

# Left: factor breakdown
add_text(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.4),
         "Top Contributing Factors", size=16, bold=True, color=CYAN)
factors = [
    ("Same audience hit 4× / 30d via SMS", 28, RED),
    ("Channel SMS — high opt-out base rate",  19, AMBER),
    ("Promotional tone in 78% of past msgs",  14, AMBER),
    ("Personal Loan saturation in segment",   12, AMBER),
    ("Region overlap with active campaigns",   9, GREEN),
]
y = Inches(1.95)
for label, pct, col in factors:
    add_text(s, Inches(0.55), y, Inches(4.4), Inches(0.4),
             label, size=11, color=LIGHT)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), y + Inches(0.42),
                            Inches(5.5), Inches(0.18))
    bg.fill.solid(); bg.fill.fore_color.rgb = DEEP; bg.line.fill.background()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), y + Inches(0.42),
                             Inches(5.5 * pct / 30), Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s, Inches(6.1), y + Inches(0.32), Inches(0.6), Inches(0.4),
             f"{pct}%", size=11, bold=True, color=col)
    y += Inches(0.78)

# Right: similar campaigns
add_text(s, Inches(7.0), Inches(1.5), Inches(6.0), Inches(0.4),
         "Top Similar Historical Campaigns (FAISS)", size=16, bold=True, color=PINK)
sim = [
    ("CMP_2241  Diwali PL Push",   "Sim 0.92", "ROI 412", GREEN),
    ("CMP_1879  Festive Loan SMS", "Sim 0.88", "ROI 287", AMBER),
    ("CMP_3104  Year-end PL Email","Sim 0.85", "ROI 521", GREEN),
    ("CMP_0917  Salary-day PL",    "Sim 0.81", "ROI 198", RED),
    ("CMP_2588  Pre-approved PL",  "Sim 0.79", "ROI 367", GREEN),
]
y = Inches(1.95)
for name, sscore, roi, col in sim:
    add_card(s, Inches(7.0), y, Inches(5.85), Inches(0.7), border=col, border_w=1.0)
    add_text(s, Inches(7.15), y + Inches(0.18), Inches(3.0), Inches(0.4),
             name, size=11, bold=True, color=LIGHT)
    add_pill(s, Inches(10.2), y + Inches(0.18), Inches(1.1), Inches(0.35),
             sscore, fill=PURPLE, color=WHITE, size=10)
    add_pill(s, Inches(11.4), y + Inches(0.18), Inches(1.3), Inches(0.35),
             roi, fill=col, color=NAVY, size=10)
    y += Inches(0.78)


# -----------------------------------------------------------------------------
# Slide 11 – Output 4: Audience Playbook + Channel Recos
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Output 4 — Audience Playbook & Channel Recommendations",
             "Actionable do / don’t guidance generated by GPT")

# Left: recommended audiences
add_card(s, Inches(0.5), Inches(1.55), Inches(6.1), Inches(2.55), border=GREEN)
add_text(s, Inches(0.7), Inches(1.65), Inches(5.7), Inches(0.4),
         "✅  Recommended Audiences", size=14, bold=True, color=GREEN)
add_bullets(s, Inches(0.7), Inches(2.1), Inches(5.7), Inches(2.0), [
    "Salaried Millennials (Tier-1) — high CTR & LTV",
    "First-Time Borrowers — high conversion intent",
    "Existing Credit Card customers — cross-sell ready",
], size=11)

# Right: avoid
add_card(s, Inches(6.7), Inches(1.55), Inches(6.1), Inches(2.55), border=RED)
add_text(s, Inches(6.9), Inches(1.65), Inches(5.7), Inches(0.4),
         "⛔  Audiences to Avoid", size=14, bold=True, color=RED)
add_bullets(s, Inches(6.9), Inches(2.1), Inches(5.7), Inches(2.0), [
    "Salaried Gen-Z on SMS — opt-out spikes",
    "Senior Citizens for Push — low engagement",
    "NRI segment — regional regulatory friction",
], size=11)

# Channels
add_card(s, Inches(0.5), Inches(4.2), Inches(6.1), Inches(2.7), border=CYAN)
add_text(s, Inches(0.7), Inches(4.3), Inches(5.7), Inches(0.4),
         "📢  Channel Recommendations", size=14, bold=True, color=CYAN)
chans = [
    ("🔥  Email",       "HIGH",   "Best ROI for HNI segment", GREEN),
    ("📌  WhatsApp",    "MEDIUM", "Good for personalised offers", AMBER),
    ("📎  SMS",         "LOW",    "Currently saturated — pause", RED),
]
y = Inches(4.75)
for icon, prio, why, col in chans:
    add_text(s, Inches(0.7), y, Inches(2.0), Inches(0.4),
             icon, size=12, bold=True, color=LIGHT)
    add_pill(s, Inches(2.6), y + Inches(0.05), Inches(1.0), Inches(0.35),
             prio, fill=col, color=NAVY, size=10)
    add_text(s, Inches(3.7), y, Inches(2.8), Inches(0.4),
             why, size=11, color=LIGHT)
    y += Inches(0.6)

# Timing
add_card(s, Inches(6.7), Inches(4.2), Inches(6.1), Inches(2.7), border=PURPLE)
add_text(s, Inches(6.9), Inches(4.3), Inches(5.7), Inches(0.4),
         "⏰  Optimal Timing", size=14, bold=True, color=PURPLE)
add_bullets(s, Inches(6.9), Inches(4.75), Inches(5.7), Inches(2.0), [
    "Best Days: Tue / Wed / Sat",
    "Best Time: 11 AM – 1 PM, 7 PM – 9 PM",
    "Avoid: Mondays, festival public holidays",
    "Recommended duration: 14–21 days",
], size=11)


# -----------------------------------------------------------------------------
# Slide 12 – Output 5: Optimizer + Executive Summary
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Output 5 — Campaign Optimizer & Executive Summary",
             "AI rewrites the campaign + produces leadership narrative")

# Optimizer before/after
add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4),
         "⚡  Campaign Optimizer (Before  →  After)", size=16, bold=True, color=CYAN)

add_card(s, Inches(0.5), Inches(2.0), Inches(6.1), Inches(2.4), border=RED)
add_text(s, Inches(0.7), Inches(2.1), Inches(5.7), Inches(0.4),
         "BEFORE — Risk: 67% (HIGH)", size=12, bold=True, color=RED)
add_text(s, Inches(0.7), Inches(2.55), Inches(5.7), Inches(1.8),
         "“Hurry! Limited time! Get instant Personal Loan with ZERO fees — apply NOW "
         "before this offer expires forever!” \nChannel: SMS  •  Audience: Salaried Millennials",
         size=11, color=LIGHT)

add_card(s, Inches(6.7), Inches(2.0), Inches(6.1), Inches(2.4), border=GREEN)
add_text(s, Inches(6.9), Inches(2.1), Inches(5.7), Inches(0.4),
         "AFTER — Risk: 31% (LOW)", size=12, bold=True, color=GREEN)
add_text(s, Inches(6.9), Inches(2.55), Inches(5.7), Inches(1.8),
         "“Pre-approved Personal Loan with waived processing fee for select Salaried "
         "customers. Check eligibility in 60 seconds.” \nChannel: Email  •  Audience: Tier-1 Salaried Millennials",
         size=11, color=LIGHT)

# Executive Summary card
add_card(s, Inches(0.5), Inches(4.55), Inches(12.3), Inches(2.35), border=PURPLE)
add_text(s, Inches(0.7), Inches(4.65), Inches(11.9), Inches(0.45),
         "🧾  Executive AI Summary", size=14, bold=True, color=PURPLE)
add_text(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.7),
         "Campaign health improved from 54 → 78 after optimization. Predicted ROI lifted "
         "from 287 to ~410 (+43%), opt-outs reduced ~50%, complaint risk halved. Channel "
         "switch from SMS → Email and a softer, eligibility-led message remove the top 3 "
         "risk factors. Recommendation: PROCEED with the optimized variant; monitor "
         "fatigue index daily for the first 7 days.",
         size=12, color=LIGHT)


# -----------------------------------------------------------------------------
# Slide 13 – Other Outputs (Live Signals, Chat, Exports)
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Additional Outputs",
             "Everything else the app surfaces during a session")

extras = [
    ("💬  Contextual AI Chat",
     "Ask follow-ups inside any section: Health, Factors, Playbook, Similar Campaigns. "
     "GPT answers using only that section’s context — no hallucinated metrics."),
    ("📡  Live Risk Signals",
     "Real-time signals like ‘fatigue spike detected’, ‘negative sentiment trending’, "
     "‘compliance keyword detected’ — surfaced from feedback & content tables."),
    ("🎯  Best Contact Time",
     "Per audience × channel, the engine returns the optimal day & hour band derived "
     "from historical engagement patterns."),
    ("🔄  Channel Alternatives",
     "Whenever channel safety < threshold, the engine suggests safer channels with "
     "improvement deltas (e.g., ‘Email +18% safety vs SMS’)."),
    ("📈  Predicted Reductions",
     "Quantified deltas vs. historical baseline: opt-outs ↓, complaints ↓, ROI ↑."),
    ("📤  Export-ready Summary",
     "JSON + executive narrative ready to paste into review decks or send to leadership."),
]
for i, (title, body) in enumerate(extras):
    cidx = i % 2
    ridx = i // 2
    x = Inches(0.5 + cidx * 6.4)
    y = Inches(1.55 + ridx * 1.75)
    add_card(s, x, y, Inches(6.2), Inches(1.55), border=CYAN)
    add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(5.8), Inches(0.4),
             title, size=13, bold=True, color=CYAN)
    add_text(s, x + Inches(0.2), y + Inches(0.55), Inches(5.8), Inches(1.0),
             body, size=11, color=LIGHT)


# -----------------------------------------------------------------------------
# Slide 14 – Business Impact
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Business Impact",
             "Measurable value for marketing operations")

kpis = [
    ("50–70%",  "Reduction in manual review effort", CYAN),
    ("₹4.3 Cr", "Estimated annual savings",          GREEN),
    ("5–10%",   "Improvement in conversion rates",   PURPLE),
    ("48K hrs", "Analyst hours freed up / year",     PINK),
]
for i, (val, lbl, col) in enumerate(kpis):
    x = Inches(0.5 + i * 3.18)
    add_card(s, x, Inches(1.7), Inches(3.0), Inches(2.0), border=col)
    add_text(s, x, Inches(1.95), Inches(3.0), Inches(0.9),
             val, size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(2.85), Inches(2.7), Inches(0.8),
             lbl, size=12, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.4),
         "Operational outcomes", size=18, bold=True, color=CYAN)
add_card(s, Inches(0.5), Inches(4.45), Inches(12.3), Inches(2.4), border=PURPLE)
add_bullets(s, Inches(0.7), Inches(4.55), Inches(12.0), Inches(2.3), [
    "Faster campaign launches — risk review collapsed from days to minutes",
    "Reduced customer fatigue — duplicate targeting flagged before send",
    "Fewer complaints & opt-outs — predictive guardrails on tone & frequency",
    "Stronger compliance — risky wording detected pre-launch",
    "Continuous learning — every new campaign enriches the FAISS index",
    "Single source of truth for campaign decisions across teams",
], size=12)


# -----------------------------------------------------------------------------
# Slide 15 – Summary / Thank You
# -----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); add_bg(s)
# Decorative bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.0), SW, Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = CYAN; bar.line.fill.background()

add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.9),
         "🛡️  OptiGuard AI in one line", size=36, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(1.2),
         "Predict marketing campaign risks BEFORE launch — with AI reasoning, "
         "historical evidence, and a one-click optimizer.",
         size=22, color=CYAN)

add_text(s, Inches(0.7), Inches(4.0), Inches(12), Inches(0.5),
         "What you get from the app", size=18, bold=True, color=PINK)
add_bullets(s, Inches(0.7), Inches(4.5), Inches(12), Inches(2.2), [
    "Campaign Health Score (0–100) + Grade + Interpretation",
    "Fatigue, Opt-out, Complaint, Trust & Channel-safety scores with reasoning",
    "Top contributing risk factors with impact %",
    "Top similar historical campaigns via FAISS semantic search",
    "Audience playbook: recommended / avoid / timing / channels",
    "AI Campaign Optimizer (rewrite) + Executive AI Summary",
], size=14)

add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         "Thank you  •  Questions?", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)

# -----------------------------------------------------------------------------
out = "OptiGuard_AI_Demo.pptx"
prs.save(out)
print(f"✅  Saved: {out}  ({len(prs.slides)} slides)")
